from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.util import Emu

from app.services.rag.parsing.base import (
    DocumentParser,
    ParsedBlock,
    ParsedDocument,
    ParsedImage,
)
from app.services.rag.parsing.image_store import ImageStore


def _shape_text(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _emu_to_int(value: Emu | int | None) -> int | None:
    if value is None:
        return None
    return int(value)


class PptxDocumentParser(DocumentParser):
    extensions = (".pptx",)

    def __init__(
        self,
        *,
        extract_images: bool = True,
        image_store: ImageStore | None = None,
    ) -> None:
        self._extract_images = extract_images
        self._image_store = image_store or ImageStore()

    def parse(self, path: str | Path) -> ParsedDocument:
        file_path = Path(path)
        try:
            presentation = Presentation(str(file_path))
            return self._to_document(presentation)
        except Exception as exc:
            raise ValueError(
                f"无法解析 PPTX 文件 {file_path.name}: {exc}"
            ) from exc

    def _to_document(self, presentation: Presentation) -> ParsedDocument:
        blocks: list[ParsedBlock] = []
        images: list[ParsedImage] = []
        slide_texts: list[str] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            heading: str | None = None
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_text = _shape_text(title_shape)
                if title_text:
                    heading = title_text

            parts: list[str] = []
            for shape in slide.shapes:
                text = _shape_text(shape)
                if text:
                    parts.append(text)

            text = "\n".join(parts)
            slide_texts.append(text)
            blocks.append(
                ParsedBlock(page=slide_number, text=text, heading=heading)
            )
            if self._extract_images:
                images.extend(
                    self._extract_slide_images(slide.shapes, slide_number=slide_number)
                )

        full_text = "\n".join(slide_texts)
        return ParsedDocument(full_text=full_text, blocks=blocks, images=images)

    def _extract_slide_images(
        self,
        shapes: object,
        *,
        slide_number: int,
    ) -> list[ParsedImage]:
        extracted: list[ParsedImage] = []
        index_in_page = 0
        for shape in shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
                blob = image.blob
                if not blob:
                    continue
                ext = (image.ext or "bin").lstrip(".").lower()
                storage_key = self._image_store.save(blob, suffix=ext)
                extracted.append(
                    ParsedImage(
                        page=slide_number,
                        storage_key=storage_key,
                        index_in_page=index_in_page,
                        width=_emu_to_int(shape.width),
                        height=_emu_to_int(shape.height),
                    )
                )
                index_in_page += 1
            except Exception:
                continue
        return extracted
