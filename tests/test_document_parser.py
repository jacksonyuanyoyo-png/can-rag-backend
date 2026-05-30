from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest

from pypdf import PdfWriter

import docx
from pptx import Presentation
from pptx.util import Inches

from app.services.rag.parsing import (
    DocxDocumentParser,
    DocumentParser,
    ParsedBlock,
    ParsedDocument,
    PdfDocumentParser,
    PptxDocumentParser,
    TxtDocumentParser,
    get_parser_for,
)
from app.services.rag.parsing.base import ParsedImage
from app.services.rag.parsing.image_store import ImageStore

MINIMAL_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
    b"\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
    b"\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01"
    b"\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _write(tmp_path: Path, name: str, data: bytes) -> Path:
    path = tmp_path / name
    path.write_bytes(data)
    return path


def test_parse_plain_text(tmp_path: Path) -> None:
    content = "第一段内容。\n\n第二段内容。\n第二段续行。"
    path = _write(tmp_path, "doc.txt", content.encode("utf-8"))

    document = TxtDocumentParser().parse(path)

    assert isinstance(document, ParsedDocument)
    assert document.full_text == content
    assert [block.text for block in document.blocks] == [
        "第一段内容。",
        "第二段内容。\n第二段续行。",
    ]


def test_parse_empty_file(tmp_path: Path) -> None:
    path = _write(tmp_path, "empty.txt", b"")

    document = TxtDocumentParser().parse(path)

    assert document.full_text == ""
    assert document.blocks == []


def test_parse_encoding_fallback(tmp_path: Path) -> None:
    path = _write(tmp_path, "bad.txt", b"good \xff bad")

    document = TxtDocumentParser().parse(path)

    assert "good" in document.full_text
    assert document.full_text == document.blocks[0].text


def test_parse_bytes_encoding_fallback() -> None:
    document = TxtDocumentParser().parse_bytes(b"\xfftext", filename="x.txt")

    assert document.full_text == "\ufffdtext"


def test_block_structure(tmp_path: Path) -> None:
    path = _write(tmp_path, "doc.txt", b"only one block")

    document = TxtDocumentParser().parse(path)

    block = document.blocks[0]
    assert isinstance(block, ParsedBlock)
    assert block.page is None
    assert block.heading is None
    assert block.text == "only one block"


def test_get_parser_for_txt() -> None:
    parser = get_parser_for("notes.TXT")

    assert isinstance(parser, TxtDocumentParser)
    assert isinstance(parser, DocumentParser)


def test_get_parser_for_text_extension() -> None:
    assert isinstance(get_parser_for("a.text"), TxtDocumentParser)


def test_get_parser_for_unsupported_returns_none() -> None:
    assert get_parser_for("report.xyz") is None


def test_get_parser_for_pdf_docx_pptx() -> None:
    assert isinstance(get_parser_for("a.pdf"), PdfDocumentParser)
    assert isinstance(get_parser_for("a.docx"), DocxDocumentParser)
    assert isinstance(get_parser_for("a.pptx"), PptxDocumentParser)


def _make_two_page_pdf(path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.add_blank_page(width=72, height=72)
    with path.open("wb") as handle:
        writer.write(handle)


def test_parse_pdf_two_pages(tmp_path: Path) -> None:
    path = tmp_path / "two.pdf"
    _make_two_page_pdf(path)

    document = PdfDocumentParser().parse(path)

    assert len(document.blocks) == 2
    assert [block.page for block in document.blocks] == [1, 2]
    assert all(block.text == "" for block in document.blocks)


def _make_sample_docx(path: Path) -> None:
    document = docx.Document()
    document.add_heading("章节标题", level=1)
    document.add_paragraph("正文第一段。")
    document.add_paragraph("正文第二段。")
    document.save(path)


def test_parse_docx_heading_and_body(tmp_path: Path) -> None:
    path = tmp_path / "sample.docx"
    _make_sample_docx(path)

    document = DocxDocumentParser().parse(path)

    assert len(document.blocks) == 1
    block = document.blocks[0]
    assert block.page is None
    assert block.heading == "章节标题"
    assert "正文第一段。" in block.text
    assert "正文第二段。" in block.text
    assert "章节标题" in document.full_text


def _make_two_slide_pptx(path: Path) -> None:
    presentation = Presentation()
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide1.shapes.title.text = "第一张标题"
    slide1.placeholders[1].text = "第一张正文"

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "第二张标题"
    slide2.placeholders[1].text = "第二张正文"

    presentation.save(path)


def test_parse_pptx_two_slides(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    _make_two_slide_pptx(path)

    document = PptxDocumentParser().parse(path)

    assert len(document.blocks) == 2
    assert [block.page for block in document.blocks] == [1, 2]
    assert document.blocks[0].heading == "第一张标题"
    assert "第一张正文" in document.blocks[0].text
    assert document.blocks[1].heading == "第二张标题"
    assert "第二张正文" in document.blocks[1].text


@pytest.mark.parametrize(
    ("parser_cls", "filename"),
    [
        (PdfDocumentParser, "bad.pdf"),
        (DocxDocumentParser, "bad.docx"),
        (PptxDocumentParser, "bad.pptx"),
    ],
)
def test_parse_corrupt_file_raises_value_error(
    tmp_path: Path,
    parser_cls: type[DocumentParser],
    filename: str,
) -> None:
    path = _write(tmp_path, filename, b"not-a-valid-file")

    with pytest.raises(ValueError, match=filename):
        parser_cls().parse(path)


def test_supports_matches_extension() -> None:
    parser = TxtDocumentParser()

    assert parser.supports("a.txt") is True
    assert parser.supports("a.docx") is False


def test_parse_bytes_not_implemented_on_base() -> None:
    class _Dummy(DocumentParser):
        extensions = (".dummy",)

        def parse(self, path):  # type: ignore[override]
            return ParsedDocument(full_text="")

    with pytest.raises(NotImplementedError):
        _Dummy().parse_bytes(b"data")


def test_image_store_save_roundtrip(tmp_path: Path) -> None:
    store = ImageStore(root=tmp_path)
    storage_key = store.save(MINIMAL_PNG, suffix="png")

    assert storage_key.startswith("kb_images/")
    assert storage_key.endswith(".png")
    saved = store.path_for(storage_key)
    assert saved.is_file()
    assert saved.read_bytes() == MINIMAL_PNG


def test_parse_pdf_without_images_has_empty_images(tmp_path: Path) -> None:
    path = tmp_path / "blank.pdf"
    _make_two_page_pdf(path)
    store = ImageStore(root=tmp_path)

    document = PdfDocumentParser(image_store=store).parse(path)

    assert document.images == []


def _make_docx_with_picture(path: Path) -> None:
    document = docx.Document()
    document.add_picture(BytesIO(MINIMAL_PNG))
    document.add_paragraph("带图正文。")
    document.save(path)


def test_parse_docx_extracts_embedded_image(tmp_path: Path) -> None:
    path = tmp_path / "with_image.docx"
    _make_docx_with_picture(path)
    store = ImageStore(root=tmp_path)

    document = DocxDocumentParser(image_store=store).parse(path)

    assert len(document.images) == 1
    image = document.images[0]
    assert isinstance(image, ParsedImage)
    assert image.page is None
    assert image.index_in_page == 0
    assert store.path_for(image.storage_key).read_bytes() == MINIMAL_PNG


def _make_pptx_with_picture(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(BytesIO(MINIMAL_PNG), Inches(1), Inches(1))
    presentation.save(path)


def test_parse_pptx_extracts_embedded_image(tmp_path: Path) -> None:
    path = tmp_path / "with_image.pptx"
    _make_pptx_with_picture(path)
    store = ImageStore(root=tmp_path)

    document = PptxDocumentParser(image_store=store).parse(path)

    assert len(document.images) == 1
    image = document.images[0]
    assert image.page == 1
    assert image.index_in_page == 0
    assert store.path_for(image.storage_key).read_bytes() == MINIMAL_PNG


def test_parse_docx_extract_images_disabled(tmp_path: Path) -> None:
    path = tmp_path / "with_image.docx"
    _make_docx_with_picture(path)
    store = ImageStore(root=tmp_path)

    document = DocxDocumentParser(
        extract_images=False, image_store=store
    ).parse(path)

    assert document.images == []
